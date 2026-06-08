#!/usr/bin/env python3
"""将 picc-hainan-report.html 的内容转换为 PowerPoint 演示文稿。"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
import os


#text第二次
# ── 颜色常量 ──────────────────────────────────────────
BG_DARK   = RGBColor(0x0F, 0x17, 0x2A)
CARD_BG   = RGBColor(0x1E, 0x29, 0x3B)
TEXT_MAIN = RGBColor(0xE2, 0xE8, 0xF0)
TEXT_DIM  = RGBColor(0x94, 0xA3, 0xB8)
RED       = RGBColor(0xE5, 0x3E, 0x3E)
ORANGE    = RGBColor(0xDD, 0x6B, 0x20)
GREEN     = RGBColor(0x38, 0xA1, 0x69)
BLUE      = RGBColor(0x31, 0x82, 0xCE)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
YELLOW    = RGBColor(0xF6, 0xAD, 0x55)

prs = Presentation()
prs.slide_width  = Inches(13.333)  # 16:9 宽屏
prs.slide_height = Inches(7.5)

# ── 工具函数 ──────────────────────────────────────────
def slide_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size=14,
                color=TEXT_MAIN, bold=False, align=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return tf

def add_rich_textbox(slide, left, top, width, height, lines, font_name='Microsoft YaHei'):
    """lines: [(text, font_size, color, bold, align), ...]"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, ln in enumerate(lines):
        text, font_size, color, bold = ln[:4]
        align = ln[4] if len(ln) > 4 else PP_ALIGN.LEFT
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = align
        p.space_after = Pt(4)
    return tf

def add_card(slide, left, top, width, height):
    """添加一个圆角卡片背景（用矩形模拟）"""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.color.rgb = RGBColor(0x33, 0x41, 0x55)
    shape.line.width = Pt(0.5)
    shape.shadow.inherit = False
    return shape

def add_accent_bar(slide, left, top, width, height, color=RED):
    shape = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_kpi_card(slide, left, top, label, value, change, change_down=False):
    add_card(slide, left, top, 2.8, 1.8)
    clr = RED if change_down else GREEN
    add_textbox(slide, left + 0.2, top + 0.15, 2.4, 0.4, label,
                font_size=10, color=TEXT_DIM, align=PP_ALIGN.CENTER)
    add_textbox(slide, left + 0.2, top + 0.5, 2.4, 0.6, value,
                font_size=24, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, left + 0.2, top + 1.15, 2.4, 0.35, change,
                font_size=10, color=clr, align=PP_ALIGN.CENTER)

# ── 封面 ──────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
slide_bg(slide)
add_accent_bar(slide, 0, 3.25, 13.333, 0.06, RED)
add_textbox(slide, 1, 1.4, 11.3, 0.8,
            '中国人保海南分公司', font_size=20, color=TEXT_DIM, align=PP_ALIGN.CENTER)
add_textbox(slide, 1, 2.0, 11.3, 1.2,
            '2025–2026 业务情况报告', font_size=42, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_textbox(slide, 1, 3.8, 11.3, 0.6,
            '数据来源：中国人保集团年报 · 南海网 · 南国都市报 · 同花顺等公开信息',
            font_size=12, color=TEXT_DIM, align=PP_ALIGN.CENTER)
add_textbox(slide, 1, 5.5, 11.3, 0.5,
            '报告日期：2026 年 6 月 7 日', font_size=11, color=TEXT_DIM, align=PP_ALIGN.CENTER)

# ── 目录 ──────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(slide)
add_accent_bar(slide, 0.8, 1.0, 0.06, 2.0, RED)
add_textbox(slide, 1.1, 0.95, 5, 0.6, '报告目录', font_size=32, color=WHITE, bold=True)

toc = [
    ('01', '集团 2025 年度核心数据'),
    ('02', '集团 2026 年一季度概况'),
    ('03', '海南分公司 2025 年度业绩'),
    ('04', '业务创新与"五向图强"战略'),
    ('05', '2026 年最新动态'),
    ('06', '2026 理赔服务升级'),
    ('07', '数据来源与说明'),
]
for i, (num, title) in enumerate(toc):
    y = 2.2 + i * 0.7
    add_textbox(slide, 1.5, y, 0.6, 0.5, num, font_size=20, color=RED, bold=True)
    add_textbox(slide, 2.2, y, 8, 0.5, title, font_size=20, color=TEXT_MAIN)

# ── 1. 集团 2025 年度核心数据 ──────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(slide)
add_accent_bar(slide, 0.6, 0.5, 0.06, 1.8, RED)
add_textbox(slide, 0.9, 0.45, 10, 0.6, '一、中国人保集团 2025 年度核心数据',
            font_size=26, color=WHITE, bold=True)
add_textbox(slide, 0.9, 1.05, 10, 0.4, '年报发布于 2026 年 3 月 26 日',
            font_size=11, color=TEXT_DIM)

kpis = [
    ('营业总收入', '6,690 亿', '↑ 7.6%', False),
    ('保险服务收入', '5,707 亿', '↑ 6.1%', False),
    ('原保费收入', '7,383 亿', '↑ 6.5%', False),
    ('归母净利润', '466 亿', '↑ 8.8%', False),
]
for i, (label, value, change, down) in enumerate(kpis):
    add_kpi_card(slide, 0.6 + i * 3.1, 1.6, label, value, change, down)

kpis2 = [
    ('总资产', '2.03 万亿', '↑ 14.8%', False),
    ('财险市场份额', '31.6%', '行业首位', False),
    ('总投资收益率', '5.7%', '—', False),
    ('人保财险 COR', '持续优化', '承保向好', False),
]
for i, (label, value, change, down) in enumerate(kpis2):
    add_kpi_card(slide, 0.6 + i * 3.1, 3.7, label, value, change, down)

# 柱状图
chart_data = CategoryChartData()
chart_data.categories = ['2023', '2024', '2025']
chart_data.add_series('营业总收入(千亿)', [5.36, 6.21, 6.69])
chart_data.add_series('保险服务收入(千亿)', [4.68, 5.38, 5.71])
chart_data.add_series('原保费收入(千亿)', [6.05, 6.93, 7.38])

chart = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(0.6), Inches(5.6), Inches(12), Inches(1.7), chart_data
).chart
chart.has_legend = True
chart.legend.include_in_layout = False
chart.legend.font.size = Pt(8)
chart.legend.font.color.rgb = TEXT_DIM
chart.chart_title.has_text_frame = False
chart.style = 2
plot = chart.plots[0]
series_colors = [RED, ORANGE, BLUE]
for idx, color in enumerate(series_colors):
    series = plot.series[idx]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = color

# ── 2. 2026 Q1 概况 ────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(slide)
add_accent_bar(slide, 0.6, 0.5, 0.06, 1.8, RED)
add_textbox(slide, 0.9, 0.45, 10, 0.6, '二、中国人保集团 2026 年一季度概况',
            font_size=26, color=WHITE, bold=True)

# 卡片
add_card(slide, 0.6, 1.5, 5.8, 3.2)
add_textbox(slide, 1.0, 1.7, 5.0, 0.5, '📊 2026 Q1 关键指标', font_size=16, color=WHITE, bold=True)

q1_data = [
    ('保险服务收入', '行业首位', '↑ 2%', True),
    ('净利润', '约 88 亿元', '↓ 31%（投资端承压）', False),
    ('承保端 COR', '持续优化', '改善', True),
    ('高价值业务占比', '不断提升', '↑', True),
    ('新业务价值 (NBV)', '稳步增长', '↑', True),
]
for i, (label, value, change, up) in enumerate(q1_data):
    y = 2.35 + i * 0.5
    add_textbox(slide, 1.0, y, 2.0, 0.4, label, font_size=12, color=TEXT_DIM)
    add_textbox(slide, 3.0, y, 1.8, 0.4, value, font_size=12, color=WHITE, bold=True)
    clr = GREEN if up else RED
    add_textbox(slide, 4.8, y, 1.5, 0.4, change, font_size=11, color=clr)

# 右侧说明
add_card(slide, 7.0, 1.5, 5.8, 3.2)
add_textbox(slide, 7.4, 1.7, 5.0, 0.5, '⚠ 关键解读', font_size=16, color=WHITE, bold=True)
add_rich_textbox(slide, 7.4, 2.3, 5.0, 2.2, [
    ('利润下滑主因：', 12, TEXT_DIM, True),
    ('一季度净利润下降 31%，主要受资本市场波动影响，投资端收益承压。', 11, TEXT_MAIN, False),
    ('', 8, TEXT_DIM, False),
    ('承保端表现向好：', 12, TEXT_DIM, True),
    ('保险服务收入仍保持正增长（+2%），综合成本率（COR）持续优化，主业盈利能力稳健。', 11, TEXT_MAIN, False),
    ('', 8, TEXT_DIM, False),
    ('结构优化：', 12, TEXT_DIM, True),
    ('高价值业务占比不断提升，新业务价值（NBV）稳步增长，业务质量持续改善。', 11, TEXT_MAIN, False),
])

add_card(slide, 0.6, 5.1, 12.1, 1.8)
add_rich_textbox(slide, 1.0, 5.3, 11.3, 1.4, [
    ('背景补充：集团层面 vs 海南分公司', 14, WHITE, True),
    ('中国人保海南分公司作为省级分支机构，不单独发布年度报告，其经营数据纳入人保财险合并报表。', 11, TEXT_DIM, False),
    ('集团年报中不披露省分公司维度的保费收入、利润等财务指标——省级分公司数据仅能通过官方新闻、媒体报道等渠道获取。', 11, TEXT_DIM, False),
    ('因此，本报告前两部分聚焦集团整体表现，第三部分起聚焦海南分公司可获取的公开数据。', 11, TEXT_DIM, False),
])

# ── 3. 海南分公司业绩 ──────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(slide)
add_accent_bar(slide, 0.6, 0.5, 0.06, 1.8, ORANGE)
add_textbox(slide, 0.9, 0.45, 10, 0.6, '三、海南分公司 2025 年度业绩概览',
            font_size=26, color=WHITE, bold=True)
add_textbox(slide, 0.9, 1.05, 10, 0.4, '数据来源：南海网、南国都市报等官方媒体报道',
            font_size=11, color=TEXT_DIM)

hn_kpis = [
    ('累计提供保障额度', '4.3 万亿元', '↑ 14%', False),
    ('累计支付赔款', '26.06 亿元', '—', False),
    ('巨灾保险年保障', '10 亿元/年', '5 年周期', False),
    ('新能源汽车承保', '9 万+ 辆', '—', False),
]
for i, (label, value, change, down) in enumerate(hn_kpis):
    add_kpi_card(slide, 0.6 + i * 3.1, 1.6, label, value, change, down)

# 趋势图
add_card(slide, 0.6, 3.7, 12.1, 3.5)
chart_data2 = CategoryChartData()
chart_data2.categories = ['2023', '2024', '2025']
chart_data2.add_series('保障额度（万亿元）', [3.0, 3.77, 4.3])
chart_data2.add_series('赔款支出（亿元）', [18.5, 22.3, 26.06])

chart2 = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(1.0), Inches(3.9), Inches(11.3), Inches(3.1), chart_data2
).chart
chart2.has_legend = True
chart2.legend.include_in_layout = False
chart2.legend.font.size = Pt(9)
chart2.legend.font.color.rgb = TEXT_DIM
chart2.chart_title.has_text_frame = False
chart2.style = 2
colors2 = [ORANGE, RGBColor(0xE5, 0x3E, 0x3E)]
for idx, color in enumerate(colors2):
    series = chart2.plots[0].series[idx]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = color

# ── 4. 五向图强 ────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(slide)
add_accent_bar(slide, 0.6, 0.5, 0.06, 1.8, GREEN)
add_textbox(slide, 0.9, 0.45, 10, 0.6, '四、业务创新与"五向图强"战略',
            font_size=26, color=WHITE, bold=True)

five = [
    ('🌾', '向种图强', '创新"粮保宝"项目\n卫星遥感 + 气象预警'),
    ('🌊', '向海图强', '全国首单"风险共担"\n海洋试验设备保险'),
    ('🚀', '向天图强', '牵头成立\n海南商业火箭发射共保体'),
    ('🌿', '向绿图强', '光伏项目 25 个\n保障额度 8.8 亿元'),
    ('💻', '向数图强', '数字资产 / 知识产权\n互联网专属产品'),
]
for i, (icon, title, desc) in enumerate(five):
    x = 0.6 + i * 2.5
    add_card(slide, x, 1.6, 2.2, 2.6)
    add_textbox(slide, x + 0.1, 1.75, 2.0, 0.5, icon, font_size=28, align=PP_ALIGN.CENTER)
    add_textbox(slide, x + 0.1, 2.3, 2.0, 0.4, title, font_size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, x + 0.1, 2.8, 2.0, 0.9, desc, font_size=10, color=TEXT_DIM, align=PP_ALIGN.CENTER)

# 其他亮点
add_card(slide, 0.6, 4.6, 12.1, 2.6)
add_textbox(slide, 1.0, 4.8, 11.0, 0.4, '🏆 其他重要创新成果', font_size=16, color=WHITE, bold=True)

innovations = [
    '巨灾保险首席承保人：牵头 14 家行业主体承保海南省巨灾保险，覆盖台风、洪水、暴雨等自然灾害，5 年周期，每年保障额度 10 亿元',
    '跨境"商旅保"：2025 年 12 月落地自贸港跨境商旅保险，为外籍人士提供保障',
    '全省首单版权侵权保险：签发全省首单版权被侵权费用损失保险',
    '自贸港金融开放样板：持续探索自贸港特色保险产品创新',
]
for i, txt in enumerate(innovations):
    y = 5.3 + i * 0.42
    add_textbox(slide, 1.0, y, 0.3, 0.35, '●', font_size=10, color=GREEN)
    add_textbox(slide, 1.3, y, 11.0, 0.35, txt, font_size=11, color=TEXT_MAIN)

# ── 5. 2026 最新动态 ────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(slide)
add_accent_bar(slide, 0.6, 0.5, 0.06, 1.8, BLUE)
add_textbox(slide, 0.9, 0.45, 10, 0.6, '五、2026 年最新动态', font_size=26, color=WHITE, bold=True)

timeline = [
    ('2026 年 2 月', '发布《农业保险承保理赔指引》，规范种植业、养殖业承保与理赔全流程'),
    ('2026 年 3 月', '集团发布 2025 年度报告：总资产突破 2 万亿，净利润增长 8.8%'),
    ('2026 年 4 月', '集团发布一季度报告：保险服务收入增 2%，投资端承压致利润降 31%'),
    ('2026 年 5 月', '海南分公司客户节启动 + 全国"理赔明白卡"上线：覆盖车险/非车险/农险'),
]

for i, (date, content) in enumerate(timeline):
    y = 1.6 + i * 1.4
    # 圆点
    dot = slide.shapes.add_shape(
        9,  # MSO_SHAPE.OVAL
        Inches(1.1), Inches(y + 0.1), Inches(0.18), Inches(0.18)
    )
    dot.fill.solid()
    dot.fill.fore_color.rgb = RED
    dot.line.fill.background()
    # 竖线
    if i < len(timeline) - 1:
        line = slide.shapes.add_shape(
            1, Inches(1.16), Inches(y + 0.35), Inches(0.04), Inches(0.9)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(0x33, 0x41, 0x55)
        line.line.fill.background()

    add_card(slide, 1.6, y, 11.0, 1.1)
    add_textbox(slide, 1.9, y + 0.1, 2.5, 0.35, date, font_size=12, color=ORANGE, bold=True)
    add_textbox(slide, 1.9, y + 0.5, 10.3, 0.5, content, font_size=13, color=TEXT_MAIN)

# ── 6. 理赔服务升级 ────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(slide)
add_accent_bar(slide, 0.6, 0.5, 0.06, 1.8, RED)
add_textbox(slide, 0.9, 0.45, 10, 0.6, '六、2026 理赔服务升级', font_size=26, color=WHITE, bold=True)

# "五亮" 服务
add_card(slide, 0.6, 1.5, 5.8, 3.0)
add_textbox(slide, 1.0, 1.7, 5, 0.4, '🔍 理赔"五亮"服务', font_size=16, color=WHITE, bold=True)
five_bright = ['亮身份：主动告知服务人员信息', '亮过程：理赔进展全程可视可溯',
               '亮指南：清晰指引每一步操作', '亮资质：展示服务团队专业能力', '亮评价：客户可实时反馈评价']
for i, txt in enumerate(five_bright):
    y = 2.25 + i * 0.42
    add_textbox(slide, 1.0, y, 0.3, 0.35, '✦', font_size=10, color=RED)
    add_textbox(slide, 1.3, y, 5.0, 0.35, txt, font_size=11, color=TEXT_MAIN)

# 理赔明白卡
add_card(slide, 7.0, 1.5, 5.8, 3.0)
add_textbox(slide, 7.4, 1.7, 5, 0.4, '📋 "理赔明白卡"服务', font_size=16, color=WHITE, bold=True)

cards_info = [
    ('车险', '事故场景理赔步骤，点击即获取'),
    ('非车险', '医疗/意外/学幼/家财/责任险'),
    ('农险', '养殖险+种植险，灾害高效理赔'),
]
for i, (name, desc) in enumerate(cards_info):
    y = 2.25 + i * 0.55
    add_textbox(slide, 7.4, y, 1.2, 0.35, name, font_size=12, color=ORANGE, bold=True)
    add_textbox(slide, 8.7, y, 3.8, 0.35, desc, font_size=11, color=TEXT_MAIN)

# 2026 客户节
add_card(slide, 0.6, 5.0, 12.1, 2.2)
add_textbox(slide, 1.0, 5.15, 11.0, 0.4, '🎉 2026 年客户节（5 月 18 日启动）', font_size=16, color=WHITE, bold=True)
add_rich_textbox(slide, 1.0, 5.6, 11.0, 1.4, [
    ('主题："倾听您的心声  服务您的需求"', 12, TEXT_MAIN, False),
    ('核心升级：推出"明白赔"功能，理赔全程可视可溯，客户可随时查看案件流程', 11, TEXT_DIM, False),
    ('同步活动：全省理赔服务体验官招募、风险减量及金融消保知识普及', 11, TEXT_DIM, False),
    ('出席领导：人保海南省分公司党委委员、副总经理朱朝晖致辞', 11, TEXT_DIM, False),
])

# ── 7. 数据来源与说明 ──────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(slide)
add_accent_bar(slide, 0.6, 0.5, 0.06, 1.8, RED)
add_textbox(slide, 0.9, 0.45, 10, 0.6, '七、数据来源与说明', font_size=26, color=WHITE, bold=True)

add_card(slide, 0.6, 1.5, 12.1, 5.5)
sources = [
    ('1', '中国人保 2026 年一季度报告 — 东方财富', 'https://finance.eastmoney.com/a/202604293724878672.html'),
    ('2', '中国人保 2025 年度报告分析 — 新浪财经', 'https://finance.sina.com.cn/stock/relnews/hk/2026-03-27/doc-inhsmamt9585424.shtml'),
    ('3', '创新保险供给 探索自贸港金融开放样板 — 南国都市报', 'http://szb.ngdsb.cn/h5/html5/2026-01/04/content_58870_19189047.htm'),
    ('4', '人保财险海南省分公司专题报道 — 南海网', 'https://v.hinews.cn/xinwen/show-1525316.html'),
    ('5', '人保财险海南省分公司 2026 年客户节 — 新海南客户端', 'https://www.hinews.cn/page?n=2824750&m=1&s=1044'),
    ('6', '人保财险推出"理赔明白卡"服务 — 同花顺', 'http://stock.10jqka.com.cn/20260518/c676760457.shtml'),
    ('7', 'PICC Q1 2026 Earnings Review — Futu', 'https://news.futunn.com/en/post/72375580'),
    ('8', 'PICC profit rises 26% — Insurance Asia', 'https://insuranceasia.com/insurance/news/picc-profit-rises-26-growth-set-ease'),
]
for i, (num, title, url) in enumerate(sources):
    y = 1.8 + i * 0.55
    add_textbox(slide, 1.0, y, 0.3, 0.35, num, font_size=10, color=RED, bold=True)
    add_textbox(slide, 1.3, y, 7.0, 0.35, title, font_size=11, color=TEXT_MAIN)
    add_textbox(slide, 1.3, y + 0.2, 10.0, 0.3, url, font_size=8, color=TEXT_DIM)

add_textbox(slide, 1.0, 6.3, 11.0, 0.5,
            '⚠ 说明：海南分公司作为省级分支机构不单独发年报，数据来自官方新闻披露。部分 2026 年数据截至 Q1。',
            font_size=10, color=YELLOW)

# ── 结尾页 ─────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(slide)
add_accent_bar(slide, 0, 3.6, 13.333, 0.06, RED)
add_textbox(slide, 1, 2.8, 11.3, 0.8, '感谢观看', font_size=42, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_textbox(slide, 1, 4.0, 11.3, 0.6, '中国人保海南分公司 2025–2026 业务情况报告',
            font_size=16, color=TEXT_DIM, align=PP_ALIGN.CENTER)
add_textbox(slide, 1, 4.7, 11.3, 0.5, '本报告基于公开信息整理，不构成投资建议',
            font_size=11, color=TEXT_DIM, align=PP_ALIGN.CENTER)

# ── 保存 ───────────────────────────────────────────────
output_path = '/Users/wugd/project/cc-project/picc-hainan-report.pptx'
prs.save(output_path)
print(f'✅ PPT 已生成: {output_path}')
print(f'   共 {len(prs.slides)} 页幻灯片')
